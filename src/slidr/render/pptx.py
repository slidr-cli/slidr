"""PPTX renderer using python-pptx. Produces native PowerPoint shapes.
Colors are extracted from the theme CSS via tinycss2."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from slidr.parser.ast import (
    Document, Heading, Paragraph, Grid, Card, Table, Quote, ListNode, AttrNode,
    Text as ASText, CodeSpan, SoftBreak,
)
from slidr.theme.colors import parse_theme, to_rgb

# Font sizes in points
FONT = {1: 44, 2: 32, 3: 18, -1: 24, 0: 18}


def _theme_colors(css: str) -> dict:
    """Extract key colors from theme CSS."""
    vars_ = parse_theme(css)
    bg = vars_.get("--bg", "#ffffff")
    ink = vars_.get("--ink", "#333333")
    muted = vars_.get("--muted", "#777777")
    accent = vars_.get("--green2", vars_.get("--accent", "#0288d1"))
    table_bg = vars_.get("--panel", "#0f1d15") if bg.startswith("#0") else "#f5f5f5"
    table_ink = vars_.get("--ink", "#333333")
    return {
        "ink": to_rgb(ink),
        "muted": to_rgb(muted),
        "accent": to_rgb(accent),
        "white": (0xFF, 0xFF, 0xFF),
        "table_bg": to_rgb(table_bg),
        "table_ink": to_rgb(table_ink),
        "bg": bg,
    }


def render(doc: Document, output_path: Path) -> None:
    dims = doc.meta.dimensions()
    sw, sh = dims[0], dims[1]
    colors = _theme_colors(doc.meta.style)

    prs = Presentation()
    prs.slide_width = Inches(sw / 96)
    prs.slide_height = Inches(sh / 96)

    slide_master = prs.slide_masters[0]
    if colors["bg"]:
        _set_bg(slide_master, colors["bg"])

    for slide in doc.slides:
        sld = prs.slides.add_slide(prs.slide_layouts[6])
        _render_slide(sld, slide, sw, sh, colors)

    prs.save(str(output_path))


def _render_slide(sld, slide, sw: int, sh: int, colors: dict):
    left = Pt(64)
    top = Pt(5)
    width = Pt(sw - 128)

    for node in slide.children:
        top = _render_node(sld, node, left, top, width, colors)


def _render_node(sld, node, left, top, width, colors):
    if isinstance(node, Heading):
        return _add_text(sld, node.content, left, top, width, node.level)
    elif isinstance(node, Paragraph):
        return _add_text(sld, node.content, left, top, width, 0)
    elif isinstance(node, Quote):
        return _add_text(sld, node.content, left, top, width, -1)
    elif isinstance(node, Table):
        return _add_table(sld, node, left, top, width)
    elif isinstance(node, Grid):
        return _add_grid(sld, node, left, top, width)
    elif isinstance(node, ListNode):
        return _add_list(sld, node, left, top, width)
    elif isinstance(node, AttrNode):
        return _add_attr(sld, node, left, top, width)
    return top


def _add_text(sld, inlines: list, left, top, width, level: int):
    sizes = {k: Pt(v) for k, v in FONT.items()}
    bold = level > 0
    size = sizes.get(level, Pt(18))

    height = Pt(24)
    txBox = sld.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    text = _render_inline(inlines)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*colors["ink"])

    lines = max(1, len(text) // 60 + 1)
    return top + Pt(18 * lines)


def _add_table(sld, table: Table, left, top, width):
    rows = len(table.rows) + 1
    cols = len(table.headers)
    height = Pt(18 * rows)

    shape = sld.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table

    for j, h in enumerate(table.headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for i, row in enumerate(table.rows):
        for j, cell_text in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = cell_text
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = RGBColor(*colors["table_ink"])

    return top + height + Pt(8)


def _add_grid(sld, grid: Grid, left, top, width):
    cols = grid.cols or len(grid.children) or 2
    gap = Pt(8)
    col_w = (width - gap * (cols - 1)) / cols
    child_top = top
    for i, child in enumerate(grid.children):
        cl = left + i * (col_w + gap)
        child_top = _render_node(sld, child, cl, top, col_w, 0, 0)
    return child_top + Pt(8)


def _add_list(sld, node: ListNode, left, top, width):
    height = Pt(14 * len(node.items))
    txBox = sld.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(node.items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(*colors["ink"])

    return top + height + Pt(4)


def _add_attr(sld, node: AttrNode, left, top, width):
    sizes = {"kicker": Pt(15), "subtitle": Pt(22), "speaker": Pt(20), "tiny": Pt(13)}  # attr font sizes
    size = sizes.get(node.type, Pt(18))

    text = node.value
    if node.type == "speaker":
        name = node.attrs.get("name", text)
        role = node.attrs.get("role", "")
        text = f"{name}\n{role}"

    height = Pt(18)
    txBox = sld.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    color = RGBColor(*colors["accent"]) if node.type == "kicker" else RGBColor(*colors["muted"])
    p.font.color.rgb = color
    p.font.bold = node.type == "kicker"

    return top + height + Pt(4)


def _render_inline(nodes: list) -> str:
    text = ""
    for n in nodes:
        if isinstance(n, ASText):
            text += n.content
        elif isinstance(n, CodeSpan):
            text += n.content
        elif isinstance(n, SoftBreak):
            text += " "
    return text


def _extract_bg(css: str) -> str:
    import re
    m = re.search(r"--bg\s*:\s*([#\w]+)", css)
    if m:
        return m.group(1)
    return ""


def _set_bg(master, color: str):
    from pptx.oxml.ns import qn
    bg = master.element.find(qn('p:cSld'))
    if bg is None:
        return
    bgPr = bg.find(qn('p:bg'))
    if bgPr is None:
        from lxml import etree
        bgPr = etree.SubElement(bg, qn('p:bg'))
    bgPr.clear()
    from lxml import etree
    solid = etree.SubElement(bgPr, qn('p:bgPr'))
    fill = etree.SubElement(solid, qn('a:solidFill'))
    clr = etree.SubElement(fill, qn('a:srgbClr'))
    clr.set('val', color.lstrip('#'))
